import sys, re, json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

MARKITDOWN_AVAILABLE = False
PPTX2MD_AVAILABLE = False

try:
    from markitdown import MarkItDown

    MARKITDOWN_AVAILABLE = True
except ImportError:
    pass

try:
    import pptx2md

    PPTX2MD_AVAILABLE = True
except ImportError:
    pass

TOOL_NAMES = {
    "pptx": "python-pptx (plain text)",
    "md": "MarkItDown (markdown)",
    "p2": "pptx2md (markdown)",
}

COLORS = {
    "pptx": {"hex": "#2563eb", "bg": "#eff6ff"},
    "md": {"hex": "#059669", "bg": "#ecfdf5"},
    "p2": {"hex": "#d97706", "bg": "#fffbeb"},
}


def extract_pptx_slides(pptx_path):
    prs = Presentation(pptx_path)
    slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts, tables, charts, links = [], [], [], []
        notes = ""
        shapes_count = 0

        sorted_shapes = sorted(
            slide.shapes, key=lambda s: (int(s.top or 0), int(s.left or 0))
        )

        for shape in sorted_shapes:
            shapes_count += 1

            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            if shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                tables.append(
                    {"rows": len(table.rows), "cols": len(table.columns), "data": rows}
                )

            if hasattr(shape, "has_chart") and shape.has_chart:
                try:
                    chart = shape.chart
                    series = []
                    for s in chart.series:
                        series.append(
                            {
                                "name": str(s.name) if s.name is not None else "",
                                "values": [
                                    round(v, 4) if isinstance(v, float) else v
                                    for v in s.values
                                ],
                            }
                        )
                    charts.append({"type": str(chart.chart_type), "series": series})
                except:
                    pass

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            links.append(
                                {"text": run.text, "url": run.hyperlink.address}
                            )

        try:
            if slide.has_notes_slide:
                t = slide.notes_slide.notes_text_frame.text.strip()
                if t:
                    notes = t
        except:
            pass

        slides.append(
            {
                "num": idx,
                "texts": texts,
                "shapes_count": shapes_count,
                "tables_count": len(tables),
                "tables": tables,
                "charts_count": len(charts),
                "charts": charts,
                "links_count": len(links),
                "links": links,
                "notes": notes,
            }
        )

    return slides


def extract_markitdown_slides(pptx_path, total):
    md = MarkItDown()
    result = md.convert(pptx_path)
    text = result.text_content
    slides = [""] * total
    pattern = r"<!-- Slide number: (\d+) -->"
    markers = list(re.finditer(pattern, text))
    for i, m in enumerate(markers):
        num = int(m.group(1))
        if 1 <= num <= total:
            start = m.end()
            end = markers[i + 1].start() if i + 1 < len(markers) else len(text)
            slides[num - 1] = text[start:end].strip()
    aligned = sum(1 for s in slides if s)
    print(f"    {aligned}/{total} slides mapped")
    return slides


def extract_pptx2md_slides(pptx_path, output_dir, total, ref_slides):
    from pptx2md.types import ConversionConfig

    out = output_dir / "_pptx2md_temp.md"
    config = ConversionConfig(
        pptx_path=Path(pptx_path),
        output_path=out,
        image_dir=None,
        disable_image=True,
        disable_color=True,
        enable_slides=True,
    )
    pptx2md.convert(config)
    text = out.read_text(encoding="utf-8")
    out.unlink(missing_ok=True)

    parts = re.split(r"\n---\n", text)
    parts = [p.strip() for p in parts]
    while parts and not parts[0]:
        parts.pop(0)
    while parts and not parts[-1]:
        parts.pop()

    if len(parts) == total:
        print(f"    {len(parts)} slides (perfect match)")
        return parts

    print(f"    {len(parts)} parts, {total} expected - aligning...")

    def _strip(s):
        return re.sub(r"\s+", "", re.sub(r"[#*_`>|]", "", s)).lower()

    ref_fps = [_strip(" ".join(r.get("texts", [])))[:100] for r in ref_slides]
    cand_fps = [_strip(p[:100]) for p in parts]
    aligned = [""] * total
    used = set()
    for ri, rf in enumerate(ref_fps):
        if not rf:
            continue
        best = (-1, 0)
        for ci, cf in enumerate(cand_fps):
            if ci in used or not cf:
                continue
            s = sum(1 for a, b in zip(rf, cf) if a == b)
            if s > best[1]:
                best = (ci, s)
        if best[0] >= 0 and best[1] >= 10:
            aligned[ri] = parts[best[0]]
            used.add(best[0])
    print(f"    {sum(1 for a in aligned if a)}/{total} aligned")
    return aligned


def strip_md_chars(text):
    text = re.sub(r"!\[.*?\]\(.*?\)", "", text)
    text = re.sub(r"#{1,6}\s*", "", text)
    text = re.sub(r"\*\*", "", text)
    text = re.sub(r"__", "", text)
    text = re.sub(r"`", "", text)
    text = re.sub(r"^\s*[-*+]\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+[.)]\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\|", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def detect_tables(text):
    return sum(1 for line in text.split("\n") if line.count("|") >= 2)


def fmt_html(text):
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
    text = re.sub(r"__(.+?)__", r"<b>\1</b>", text)
    text = re.sub(r"`([^`]+)`", r"<code>\1</code>", text)
    text = re.sub(r"^# (.+)$", r'<span class="h1">\1</span>', text, flags=re.MULTILINE)
    text = re.sub(r"^## (.+)$", r'<span class="h2">\1</span>', text, flags=re.MULTILINE)
    text = re.sub(
        r"^### (.+)$", r'<span class="h3">\1</span>', text, flags=re.MULTILINE
    )
    text = re.sub(r"!\[.*?\]\(.*?\)", "[image]", text)
    return text


def compute_scores(all_data, total):
    tools = ["pptx", "md", "p2"]
    results = {}
    for key in tools:
        total_content = sum(d[f"{key}_chars"] for d in all_data)
        total_raw = sum(d[f"{key}_raw"] for d in all_data)
        empty = sum(1 for d in all_data if not d[f"{key}_text"])
        tables = sum(d[f"{key}_tables"] for d in all_data)
        results[key] = {
            "content_chars": total_content,
            "raw_chars": total_raw,
            "noise_chars": total_raw - total_content,
            "empty_slides": empty,
            "tables": tables,
        }

    max_content = max(r["content_chars"] for r in results.values()) or 1
    for key, r in results.items():
        content_score = r["content_chars"] / max_content * 100
        cleanliness = (
            (r["content_chars"] / r["raw_chars"] * 100) if r["raw_chars"] else 0
        )
        coverage = (total - r["empty_slides"]) / total * 100
        r["content_score"] = round(content_score, 1)
        r["cleanliness"] = round(cleanliness, 1)
        r["coverage"] = round(coverage, 1)
        r["overall"] = round(
            0.35 * content_score + 0.30 * cleanliness + 0.35 * coverage, 1
        )

    return results


def build_html(pptx_path, pptx_slides, md_slides, p2md_slides, total, output_dir):
    data = []
    for i in range(total):
        p = pptx_slides[i]
        md = md_slides[i] if md_slides and i < len(md_slides) else ""
        p2 = p2md_slides[i] if p2md_slides and i < len(p2md_slides) else ""

        pptx_text = "\n".join(p["texts"])
        data.append(
            {
                "num": p["num"],
                "pptx_text": pptx_text,
                "md_text": md,
                "p2_text": p2,
                "shapes": p["shapes_count"],
                "tables": p["tables_count"],
                "notes": p["notes"],
                "pptx_chars": len(strip_md_chars(pptx_text)),
                "pptx_raw": len(pptx_text),
                "pptx_tables": detect_tables(pptx_text),
                "md_chars": len(strip_md_chars(md)),
                "md_raw": len(md),
                "md_tables": detect_tables(md),
                "p2_chars": len(strip_md_chars(p2)),
                "p2_raw": len(p2),
                "p2_tables": detect_tables(p2),
            }
        )

    scores = compute_scores(data, total)
    fname = Path(pptx_path).name
    c = COLORS
    json_data = json.dumps(data)
    t = TOOL_NAMES

    def score_card(key):
        s = scores[key]
        bar_pct = min(s["overall"], 100)
        return f"""
        <div class="score-card" style="border-left:4px solid {c[key]["hex"]}">
          <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px">
            <span style="font-weight:600;font-size:15px;color:{c[key]["hex"]}">{t[key]}</span>
            <span style="font-size:28px;font-weight:700;color:{c[key]["hex"]}">{s["overall"]}%</span>
          </div>
          <div class="bar-track"><div class="bar-fill" style="width:{bar_pct}%;background:{c[key]["hex"]}"></div></div>
          <div class="score-metrics">
            <span>Content: <b>{s["content_score"]}%</b></span>
            <span>Cleanliness: <b>{s["cleanliness"]}%</b></span>
            <span>Coverage: <b>{s["coverage"]}%</b></span>
          </div>
          <div class="score-detail">{s["content_chars"]} content chars &middot; {s["noise_chars"]} markup stripped &middot; {s["empty_slides"]} empty &middot; {s["tables"]} tables</div>
        </div>"""

    avails = []
    if MARKITDOWN_AVAILABLE:
        avails.append("MarkItDown")
    if PPTX2MD_AVAILABLE:
        avails.append("pptx2md")

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>PPTX Report - {fname}</title>
<style>
  *{{margin:0;padding:0;box-sizing:border-box}}
  body{{font-family:system-ui,-apple-system,'Segoe UI',Roboto,sans-serif;background:#f0f2f5;color:#1a1a2e}}
  .header{{background:linear-gradient(135deg,#1a1a2e,#16213e);color:white;padding:28px 40px}}
  .header h1{{font-size:22px;font-weight:600;margin-bottom:4px}}
  .header .sub{{opacity:.8;font-size:13px}}
  .bar-track{{height:8px;background:#e5e7eb;border-radius:4px;margin:8px 0}}
  .bar-fill{{height:8px;border-radius:4px;transition:width .3s}}
  .score-grid{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:16px;padding:0 40px 24px}}
  .score-card{{background:white;border-radius:12px;padding:20px;box-shadow:0 1px 4px rgba(0,0,0,.08)}}
  .score-metrics{{display:flex;gap:16px;font-size:12px;color:#6b7280;margin:8px 0}}
  .score-metrics b{{color:#1a1a2e}}
  .score-detail{{font-size:12px;color:#6b7280;line-height:1.6;margin-top:4px}}
  .score-label{{font-size:11px;text-transform:uppercase;letter-spacing:.5px;color:#6b7280;margin-bottom:2px}}
  .formulas{{background:white;border-radius:12px;padding:24px;margin:20px 40px;box-shadow:0 1px 4px rgba(0,0,0,.08)}}
  .formulas h3{{font-size:15px;margin-bottom:12px;color:#374151}}
  .formula-grid{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:16px}}
  .formula-item{{padding:12px;background:#f9fafb;border-radius:8px}}
  .formula-item .name{{font-weight:600;font-size:13px;color:#1a1a2e}}
  .formula-item .formula{{font-family:monospace;font-size:12px;color:#6b7280;margin:4px 0}}
  .formula-item .desc{{font-size:12px;color:#6b7280;line-height:1.5}}
  .tabs{{display:flex;gap:0;padding:0 40px;background:white;border-bottom:1px solid #e5e7eb;position:sticky;top:0;z-index:10}}
  .tabs button{{padding:12px 24px;border:none;background:none;cursor:pointer;font-size:13px;font-weight:500;color:#6b7280;border-bottom:3px solid transparent;transition:all .15s}}
  .tabs button:hover{{color:#1a1a2e}}
  .tabs button.active{{color:#2563eb;border-bottom-color:#2563eb}}
  .tab-content{{display:none}}
  .tab-content.active{{display:block}}
  .nav-bar{{display:flex;align-items:center;gap:16px;padding:12px 40px;background:white;border-bottom:1px solid #e5e7eb;flex-wrap:wrap}}
  .nav-bar button{{background:#2563eb;color:white;border:none;padding:8px 18px;border-radius:6px;cursor:pointer;font-size:14px;font-weight:500}}
  .nav-bar button:hover{{background:#1d4ed8}}
  .nav-bar button:disabled{{background:#93c5fd;cursor:default}}
  .nav-bar .slide-indicator{{font-size:15px;font-weight:600;min-width:110px;text-align:center}}
  .content{{padding:20px 40px}}
  .cols{{display:grid;grid-template-columns:1fr 1fr 1fr;gap:16px}}
  .col{{background:white;border-radius:10px;box-shadow:0 1px 4px rgba(0,0,0,.08);overflow:hidden}}
  .col-header{{padding:10px 16px;font-weight:600;font-size:13px;text-transform:uppercase;letter-spacing:.3px;border-bottom:3px solid;display:flex;justify-content:space-between;align-items:center}}
  .col-header.pptx{{color:{c["pptx"]["hex"]};border-color:{c["pptx"]["hex"]};background:{c["pptx"]["bg"]}}}
  .col-header.md{{color:{c["md"]["hex"]};border-color:{c["md"]["hex"]};background:{c["md"]["bg"]}}}
  .col-header.p2{{color:{c["p2"]["hex"]};border-color:{c["p2"]["hex"]};background:{c["p2"]["bg"]}}}
  .slide-badge{{font-size:11px;padding:2px 8px;border-radius:10px;font-weight:600}}
  .slide-badge.ok{{background:#d1fae5;color:#065f46}}
  .col-body{{padding:16px;font-size:13px;line-height:1.6;white-space:pre-wrap;word-break:break-word;font-family:'JetBrains Mono','Fira Code',monospace;min-height:250px;max-height:65vh;overflow-y:auto}}
  .col-body .empty{{color:#9ca3af;font-style:italic}}
  .col-body .h1{{display:block;font-size:16px;font-weight:700;margin:8px 0 4px;color:#1a1a2e}}
  .col-body .h2{{display:block;font-size:14px;font-weight:600;margin:6px 0 3px;color:#2d3748}}
  .col-body .h3{{display:block;font-size:13px;font-weight:600;margin:4px 0 2px;color:#4a5568}}
  .col-body code{{background:#f3f4f6;padding:1px 4px;border-radius:3px;font-size:12px}}
  .col-body b{{color:#1a1a2e}}
  .slide-stats{{display:flex;gap:8px;padding:0 40px 20px;flex-wrap:wrap}}
  .slide-stats .stat{{background:white;border-radius:8px;padding:6px 14px;font-size:12px;box-shadow:0 1px 3px rgba(0,0,0,.06)}}
  .slide-stats .stat .num{{font-weight:700;font-size:15px}}
  .notes-banner{{background:#fef9c3;border:1px solid #eab308;border-radius:6px;padding:8px 12px;margin:8px 0 0;font-size:12px}}
  .notes-banner summary{{cursor:pointer;font-weight:600;color:#854d0e}}
  .notes-banner .notes-content{{margin-top:6px;white-space:pre-wrap;font-family:monospace;font-size:12px;color:#713f12}}
  .footer{{text-align:center;padding:24px;color:#9ca3af;font-size:12px}}
  .warning{{background:#fef3c7;border:1px solid #f59e0b;border-radius:8px;padding:10px 20px;margin:16px 40px 0;font-size:13px;color:#92400e;display:flex;align-items:center;gap:8px}}
  .warning a{{color:#78350f}}
  @media(max-width:1100px){{.score-grid,.formula-grid{{grid-template-columns:1fr}}.cols{{grid-template-columns:1fr}}.header,.tabs,.nav-bar,.content{{padding-left:16px;padding-right:16px}}}}
</style>
</head>
<body>
<div class="header">
  <h1>PPTX Extraction Report</h1>
  <div class="sub">{fname} &middot; {total} slides</div>
</div>

<div class="tabs">
  <button class="active" onclick="switchTab('scores')">&#128202; Scores</button>
  <button onclick="switchTab('slides')">&#128214; Slide-by-Slide</button>
</div>

<div class="tab-content active" id="tab-scores">
  <div class="formulas">
    <h3>How Scores Are Calculated</h3>
    <div class="formula-grid">
      <div class="formula-item">
        <div class="name">Content <span style="font-weight:400;color:#6b7280">(35%)</span></div>
        <div class="formula">tool_chars / max_chars &times; 100</div>
        <div class="desc">Volume of text relative to best tool. Markdown stripped before counting.</div>
      </div>
      <div class="formula-item">
        <div class="name">Cleanliness <span style="font-weight:400;color:#6b7280">(30%)</span></div>
        <div class="formula">content_chars / raw_chars &times; 100</div>
        <div class="desc">Signal vs noise. python-pptx (plain text) scores near 100%; markdown tools lower.</div>
      </div>
      <div class="formula-item">
        <div class="name">Coverage <span style="font-weight:400;color:#6b7280">(35%)</span></div>
        <div class="formula">(total - empty) / total &times; 100</div>
        <div class="desc">Percentage of slides with extracted text.</div>
      </div>
    </div>
  </div>

  <div class="score-grid">
    {score_card("pptx")}
    {score_card("md") if MARKITDOWN_AVAILABLE else '<div class="score-card" style="border-left:4px solid #9ca3af;opacity:.6"><div style="font-weight:600;font-size:15px;color:#9ca3af">MarkItDown</div><div class="score-detail" style="margin-top:8px">Not installed &middot; pip install markitdown</div></div>'}
    {score_card("p2") if PPTX2MD_AVAILABLE else '<div class="score-card" style="border-left:4px solid #9ca3af;opacity:.6"><div style="font-weight:600;font-size:15px;color:#9ca3af">pptx2md</div><div class="score-detail" style="margin-top:8px">Not installed &middot; pip install pptx2md</div></div>'}
  </div>
</div>

<div class="tab-content" id="tab-slides">
  <div class="nav-bar">
    <button id="prevBtn" onclick="navigate(-1)">&#9664; Prev</button>
    <span class="slide-indicator" id="slideIndicator">Slide 1 / {total}</span>
    <button id="nextBtn" onclick="navigate(1)">Next &#9654;</button>
  </div>
  <div class="slide-stats" id="slideStats"></div>
  <div class="content">
    <div class="cols">
      <div class="col">
        <div class="col-header pptx">python-pptx <span class="slide-badge ok" id="pptxBadge">slide #</span></div>
        <div class="col-body" id="pptxBody"></div>
      </div>
      <div class="col">
        <div class="col-header md">MarkItDown <span class="slide-badge ok" id="mdBadge">slide #</span></div>
        <div class="col-body" id="mdBody"></div>
      </div>
      <div class="col">
        <div class="col-header p2">pptx2md <span class="slide-badge ok" id="p2Badge">slide #</span></div>
        <div class="col-body" id="p2Body"></div>
      </div>
    </div>
    <div id="notesSection" style="margin-top:12px"></div>
  </div>
</div>

<div class="footer">Generated by pptx-report &middot; {datetime.now().strftime("%Y-%m-%d %H:%M")}</div>

<script>
const DATA = {json_data};
let current = 0;

function switchTab(name){{
  document.querySelectorAll('.tab-content').forEach(t=>t.classList.remove('active'));
  document.querySelectorAll('.tabs button').forEach(b=>b.classList.remove('active'));
  document.getElementById('tab-'+name).classList.add('active');
  event.target.classList.add('active');
}}

function render(slide){{
  const d=DATA[slide];
  document.getElementById('slideIndicator').textContent='Slide '+d.num+' / '+DATA.length;
  document.getElementById('prevBtn').disabled=slide===0;
  document.getElementById('nextBtn').disabled=slide===DATA.length-1;
  ['pptx','md','p2'].forEach(function(key){{
    var body=document.getElementById(key+'Body');
    var txt=d[key+'_text'];
    if(txt&&txt.trim()){{body.innerHTML=fmt(txt)}}else{{body.innerHTML='<span class=\\"empty\\">(empty)</span>'}}
    document.getElementById(key+'Badge').textContent='slide '+d.num;
  }});
  var h='<div class=\\"slide-stats\\">'+
    '<div class=\\"stat\\">Shapes <span class=\\"num\\">'+d.shapes+'</span></div>'+
    '<div class=\\"stat\\">Tables <span class=\\"num\\">'+d.tables+'</span></div>'+
    '<div class=\\"stat\\">pptx <span class=\\"num\\">'+d.pptx_chars+'</span>c</div>'+
    '<div class=\\"stat\\">MD <span class=\\"num\\">'+d.md_chars+'</span>c</div>'+
    '<div class=\\"stat\\">p2 <span class=\\"num\\">'+d.p2_chars+'</span>c</div>'+
    '</div>';
  document.getElementById('slideStats').innerHTML=h;
  if(d.notes){{
    document.getElementById('notesSection').innerHTML='<details class=\\"notes-banner\\"><summary>&#128172; Notes (Slide '+d.num+')</summary><div class=\\"notes-content\\">'+esc(d.notes)+'</div></details>';
  }}else{{document.getElementById('notesSection').innerHTML=''}}
}}

function navigate(dir){{
  var n=current+dir;
  if(n<0||n>=DATA.length)return;
  current=n;render(current);
}}

function esc(s){{return s.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;')}}

function fmt(text){{
  return esc(text)
    .replace(/\\*\\*(.+?)\\*\\*/g,'<b>$1</b>')
    .replace(/__(.+?)__/g,'<b>$1</b>')
    .replace(/`([^`]+)`/g,'<code>$1</code>')
    .replace(/^# (.+)$/gm,'<span class=\\"h1\\">$1</span>')
    .replace(/^## (.+)$/gm,'<span class=\\"h2\\">$1</span>')
    .replace(/^### (.+)$/gm,'<span class=\\"h3\\">$1</span>')
    .replace(/!\\[.*?\\]\\(.*?\\)/g,'[image]');
}}

document.addEventListener('keydown',function(e){{if(e.key==='ArrowLeft')navigate(-1);if(e.key==='ArrowRight')navigate(1)}});
render(0);
</script>
</body>
</html>"""


def build_markdown(pptx_path, pptx_slides, md_slides, p2md_slides, total, pptx_scores):
    lines = []
    fname = Path(pptx_path).stem
    total_chars_pptx = 0
    total_tables = 0
    total_charts = 0
    total_links = 0
    total_notes = 0
    empty = 0
    per_slide = []

    lines.append(f"# {fname}")
    lines.append(f"")
    lines.append(
        f"> Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} &middot; {total} slides"
    )
    lines.append(f"")
    lines.append(f"---")
    lines.append(f"")

    for i in range(total):
        p = pptx_slides[i]
        md = md_slides[i] if md_slides and i < len(md_slides) else ""
        p2 = p2md_slides[i] if p2md_slides and i < len(p2md_slides) else ""
        slide_chars = 0
        h1 = bool(p["texts"] or p["tables"] or p["charts"] or p["links"])

        lines.append(f"## Slide {p['num']}")
        lines.append(f"")

        if not h1 and not p["notes"]:
            lines.append(f"*empty*")
            lines.append(f"")
            empty += 1
            per_slide.append(0)
            lines.append(f"---")
            lines.append(f"")
            continue

        for t in p["texts"]:
            lines.append(f"{t}")
            lines.append(f"")
            slide_chars += len(t)

        for t in p["tables"]:
            total_tables += 1
            h = t["data"][0]
            pad = lambda r: r + [""] * (len(h) - len(r))
            lines.append(f"| {' | '.join(h)} |")
            lines.append(f"| {' | '.join(['---'] * len(h))} |")
            for r in t["data"][1:]:
                lines.append(f"| {' | '.join(pad(r))} |")
            lines.append(f"")

        for c in p["charts"]:
            total_charts += 1
            ct = re.sub(r"\((\d+)\)", "", c["type"]).strip()
            lines.append(f"**Chart:** {ct}")
            for sr in c["series"]:
                vals = ", ".join(str(v) for v in sr["values"][:12])
                lbl = sr["name"] or "Series"
                lines.append(f"- {lbl}: {vals}")
                if len(sr["values"]) > 12:
                    lines.append(f"  *(+{len(sr['values']) - 12} more)*")
            lines.append(f"")

        for lk in p["links"]:
            total_links += 1
            lines.append(f"- [{lk['text']}]({lk['url']})")
        if p["links"]:
            lines.append(f"")

        if md:
            md_stripped = strip_md_chars(md)
        else:
            md_stripped = ""
        if p2:
            p2_stripped = strip_md_chars(p2)
        else:
            p2_stripped = ""

        if p["notes"]:
            total_notes += 1
            lines.append(f"> **Notes:** {p['notes']}")
            lines.append(f"")

        per_slide.append(slide_chars)
        total_chars_pptx += slide_chars
        lines.append(f"---")
        lines.append(f"")

    # Evaluation
    cov = round((total - empty) / total * 100, 1) if total else 0
    avg = round(total_chars_pptx / total, 1) if total else 0
    mx = max(per_slide) if per_slide else 0
    mn = min((c for c in per_slide if c > 0), default=0)
    std = (
        round((sum((c - avg) ** 2 for c in per_slide) / total) ** 0.5, 1)
        if total
        else 0
    )

    density = (
        min(100, round(total_chars_pptx / max(1, total * 50) * 100, 1)) if total else 0
    )
    richness_val = total_tables + total_charts + total_links + total_notes
    richness_score = (
        min(100, round(richness_val / max(1, total) * 50, 1)) if total else 0
    )
    overall = round(0.5 * cov + 0.3 * density + 0.2 * richness_score, 1)

    lines.append(f"## Evaluation")
    lines.append(f"")
    lines.append(f"### Extraction Stats (python-pptx)")
    lines.append(f"")
    lines.append(f"| Metric | Value |")
    lines.append(f"|--------|-------|")
    lines.append(f"| Total Slides | {total} |")
    lines.append(f"| Total Characters | {total_chars_pptx} |")
    lines.append(f"| Avg per Slide | {avg} |")
    lines.append(f"| Max per Slide | {mx} |")
    lines.append(f"| Min (non-empty) | {mn} |")
    lines.append(f"| Std Deviation | {std} |")
    lines.append(f"| Tables | {total_tables} |")
    lines.append(f"| Charts | {total_charts} |")
    lines.append(f"| Hyperlinks | {total_links} |")
    lines.append(f"| Slides with Notes | {total_notes} |")
    lines.append(f"| Empty Slides | {empty} |")
    lines.append(f"")
    lines.append(f"### Tool Comparison Scores")
    lines.append(f"")
    lines.append(f"| Score | python-pptx | MarkItDown | pptx2md |")
    lines.append(f"|-------|-------------|------------|---------|")

    md_scores = {}
    p2_scores = {}
    if md_slides:
        md_tc = sum(len(strip_md_chars(s)) for s in md_slides)
        md_tr = sum(len(s) for s in md_slides)
        md_em = sum(1 for s in md_slides if not s.strip())
        mc = max(md_tc, 1)
        md_cs = md_tc / mc * 100
        md_cl = (md_tc / md_tr * 100) if md_tr else 0
        md_cv = (total - md_em) / total * 100
        md_scores = {
            "content": round(md_cs, 1),
            "clean": round(md_cl, 1),
            "cov": round(md_cv, 1),
            "overall": round(0.35 * md_cs + 0.30 * md_cl + 0.35 * md_cv, 1),
        }
    if p2md_slides:
        p2_tc = sum(len(strip_md_chars(s)) for s in p2md_slides)
        p2_tr = sum(len(s) for s in p2md_slides)
        p2_em = sum(1 for s in p2md_slides if not s.strip())
        p2c = max(p2_tc, 1)
        p2_cs = p2_tc / p2c * 100
        p2_cl = (p2_tc / p2_tr * 100) if p2_tr else 0
        p2_cv = (total - p2_em) / total * 100
        p2_scores = {
            "content": round(p2_cs, 1),
            "clean": round(p2_cl, 1),
            "cov": round(p2_cv, 1),
            "overall": round(0.35 * p2_cs + 0.30 * p2_cl + 0.35 * p2_cv, 1),
        }

    lines.append(
        f"| Content (35%) | {pptx_scores['content_score']}% | {md_scores.get('content', 'N/A')}% | {p2_scores.get('content', 'N/A')}% |"
    )
    lines.append(
        f"| Cleanliness (30%) | {pptx_scores['cleanliness']}% | {md_scores.get('clean', 'N/A')}% | {p2_scores.get('clean', 'N/A')}% |"
    )
    lines.append(
        f"| Coverage (35%) | {pptx_scores['coverage']}% | {md_scores.get('cov', 'N/A')}% | {p2_scores.get('cov', 'N/A')}% |"
    )
    lines.append(
        f"| **Overall** | **{pptx_scores['overall']}%** | **{md_scores.get('overall', 'N/A')}%** | **{p2_scores.get('overall', 'N/A')}%** |"
    )
    lines.append(f"")
    lines.append(f"### Quality Score (python-pptx)")
    lines.append(f"")
    lines.append(f"| Metric | Weight | Value |")
    lines.append(f"|--------|--------|-------|")
    lines.append(f"| Coverage | 50% | {cov}% |")
    lines.append(f"| Text Density | 30% | {density}% |")
    lines.append(f"| Richness | 20% | {richness_score}% |")
    lines.append(f"| **Overall** | **100%** | **{overall}%** |")
    lines.append(f"")
    lines.append(f"---")
    lines.append(f"*Generated by pptx-report*")
    lines.append(f"")

    return "\n".join(lines)


def process_file(pptx_file, output_dir):
    fname = Path(pptx_file).stem
    file_output_dir = output_dir / fname
    file_output_dir.mkdir(parents=True, exist_ok=True)

    print(f"\n{'=' * 60}")
    print(f"Processing: {Path(pptx_file).name}")
    print(f"{'=' * 60}")

    print("[1] Extracting with python-pptx...")
    pptx_slides = extract_pptx_slides(pptx_file)
    total = len(pptx_slides)
    total_notes = sum(1 for s in pptx_slides if s["notes"])
    total_shapes = sum(s["shapes_count"] for s in pptx_slides)
    total_tables = sum(s["tables_count"] for s in pptx_slides)
    total_charts = sum(s["charts_count"] for s in pptx_slides)
    total_links = sum(s["links_count"] for s in pptx_slides)
    print(
        f"    {total} slides, {total_shapes} shapes, {total_tables} tables, {total_charts} charts"
    )
    print(f"    {total_links} links, {total_notes} slides with notes")

    md_slides = None
    p2md_slides = None

    if MARKITDOWN_AVAILABLE:
        print("[2] Extracting with MarkItDown...")
        md_slides = extract_markitdown_slides(pptx_file, total)
    else:
        print("[2] MarkItDown not available (pip install markitdown)")

    if PPTX2MD_AVAILABLE:
        print("[3] Extracting with pptx2md...")
        p2md_slides = extract_pptx2md_slides(
            pptx_file, file_output_dir, total, pptx_slides
        )
    else:
        print("[3] pptx2md not available (pip install pptx2md)")

    print("\n[4] Generating HTML report...")
    html = build_html(
        pptx_file, pptx_slides, md_slides, p2md_slides, total, file_output_dir
    )
    report_path = file_output_dir / "report.html"
    report_path.write_text(html, encoding="utf-8")
    print(f"    Report: {report_path.resolve()}")

    # Compute scores for all tools
    score_data = [
        {
            "pptx_text": "\n".join(s["texts"]),
            "pptx_chars": len(strip_md_chars("\n".join(s["texts"]))),
            "pptx_raw": len("\n".join(s["texts"])),
            "pptx_tables": s["tables_count"],
            "md_text": md_slides[i] if md_slides else "",
            "md_chars": len(strip_md_chars(md_slides[i] if md_slides else "")),
            "md_raw": len(md_slides[i] if md_slides else ""),
            "md_tables": detect_tables(md_slides[i] if md_slides else ""),
            "p2_text": p2md_slides[i] if p2md_slides else "",
            "p2_chars": len(strip_md_chars(p2md_slides[i] if p2md_slides else "")),
            "p2_raw": len(p2md_slides[i] if p2md_slides else ""),
            "p2_tables": detect_tables(p2md_slides[i] if p2md_slides else ""),
        }
        for i, s in enumerate(pptx_slides)
    ]
    all_scores = compute_scores(score_data, total)
    print(
        f"    python-pptx: {all_scores['pptx']['overall']}% | MarkItDown: {all_scores['md']['overall']}% | pptx2md: {all_scores['p2']['overall']}%"
    )

    print("\n[5] Generating Markdown (with evaluation)...")
    markdown = build_markdown(
        pptx_file, pptx_slides, md_slides, p2md_slides, total, all_scores["pptx"]
    )
    md_path = file_output_dir / "report.md"
    md_path.write_text(markdown, encoding="utf-8")
    print(f"    Markdown: {md_path.resolve()} ({len(markdown)} chars)")

    return {
        "file": Path(pptx_file).name,
        "slides": total,
        "shapes": total_shapes,
        "tables": total_tables,
        "charts": total_charts,
        "links": total_links,
        "notes": total_notes,
        "html": str(report_path),
        "md": str(md_path),
    }


def main():
    import argparse

    parser = argparse.ArgumentParser(
        description="PPTX to HTML + Markdown report with quality evaluation"
    )
    parser.add_argument("input", help="Path to .pptx file")
    parser.add_argument(
        "-o", "--output", default="outputs", help="Output directory (default: outputs/)"
    )
    parser.add_argument(
        "--no-md", action="store_true", help="Skip MarkItDown extraction"
    )
    parser.add_argument("--no-p2", action="store_true", help="Skip pptx2md extraction")
    args = parser.parse_args()

    path = Path(args.input)
    if not path.exists():
        print(f"Error: '{path}' not found", file=sys.stderr)
        sys.exit(1)

    if args.no_md:
        global MARKITDOWN_AVAILABLE
        MARKITDOWN_AVAILABLE = False
    if args.no_p2:
        global PPTX2MD_AVAILABLE
        PPTX2MD_AVAILABLE = False

    output_dir = Path(args.output)
    output_dir.mkdir(exist_ok=True)

    try:
        r = process_file(str(path.resolve()), output_dir)
        print(f"\n{'=' * 60}")
        print(f"Done: {r['file']}")
        print(f"  HTML:  {r['html']}")
        print(f"  MD:    {r['md']}")
        print(f"{'=' * 60}")
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
